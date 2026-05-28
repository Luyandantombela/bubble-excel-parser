import io, json, os, re, datetime
import pandas as pd
from flask import Flask, jsonify, request, send_file
from flask_cors import CORS

# ── New imports for /extract-to-excel ──
try:
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

try:
    from spellchecker import SpellChecker
    _SPELL = SpellChecker()
    SPELLCHECK_AVAILABLE = True
except ImportError:
    _SPELL = None
    SPELLCHECK_AVAILABLE = False

app = Flask(__name__)
CORS(app)

CUSTOM_WHITELIST = {
    'joburg', 'pretoria', 'durban', 'bongi', 'tina', 'chris', 'linda',
    'sarah', 'mike', 'amy', 'john', 'todo', 'asdf', 'idk'
}

# ── Pre-compile every regex used in hot loops ──
_RE_NON_ALPHA     = re.compile(r'[^a-zA-Z\s]')
_RE_NON_DIGIT_SEP = re.compile(r'[^\d\.\-]')
_RE_DATE1         = re.compile(r'^\d{4}[-/]\d{2}[-/]\d{2}$')
_RE_DATE2         = re.compile(r'^\d{2}[-/]\d{2}[-/]\d{2,4}$')
_RE_DATE_MASK     = re.compile(r'\d')
_RE_VALID_EMAIL   = re.compile(r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,4}$')
_RE_NON_DIGIT     = re.compile(r'[\s\-\(\)\+]')
_RE_NUMBER_CLEAN  = re.compile(r'[^\d\.,\-]')
_RE_CONTAMINATE   = re.compile(r'[A-Za-z\$£€R]')
_RE_TITLE         = re.compile(r'^[A-Z][a-z]*(\s+[A-Z][a-z]*)*$')
_RE_NON_ALPHA_KEY = re.compile(r'[^a-zA-Z]')

_DATE_FORMATS = [
    '%Y/%m/%d', '%Y-%m-%d', '%d/%m/%Y', '%d-%m-%Y',
    '%m-%d-%y', '%m/%d/%Y', '%d %b %Y', '%b %d, %Y', '%m-%d-%Y'
]

_TEXT_NUMBERS = {
    'one','two','three','four','five','six','seven','eight','nine','ten',
    'twenty','thirty','forty','fifty'
}


# ─────────────────────────────────────────────
# SPELL CHECK
# ─────────────────────────────────────────────
def find_column_typos(series_str):
    if SPELLCHECK_AVAILABLE:
        found = {}
        for val in series_str:
            clean_text = _RE_NON_ALPHA.sub('', val)
            words = [w.lower() for w in clean_text.split() if len(w) > 2]
            unknown = _SPELL.unknown(words)
            for word in unknown:
                if word not in CUSTOM_WHITELIST and word not in found:
                    correction = _SPELL.correction(word)
                    found[word] = correction if correction and correction != word else None
        return [{"original": k, "suggestion": v} for k, v in found.items()]

    word_pool = []
    for val in series_str:
        clean_text = _RE_NON_ALPHA.sub('', val)
        for chunk in clean_text.split():
            if len(chunk) > 3:
                word_pool.append(chunk.lower())
    if not word_pool:
        return []

    freq_map = pd.Series(word_pool).value_counts()
    unique_words = list(freq_map.head(200).index)
    flagged = set()
    for i, word in enumerate(unique_words):
        for candidate in unique_words[i + 1:]:
            if abs(len(word) - len(candidate)) <= 2:
                if (word[:-1] == candidate or candidate[:-1] == word or
                        (word[:-2] == candidate[:-2] and word[:3] == candidate[:3])):
                    if not (word.startswith("unite") or candidate.startswith("unite")):
                        rare = word if freq_map[word] < freq_map[candidate] else candidate
                        flagged.add(rare)
    return [{"original": w, "suggestion": None} for w in flagged]


# ─────────────────────────────────────────────
# TYPE DETECTION
# ─────────────────────────────────────────────
def analyze_exclusive_type(values: list):
    filled_count = len(values)
    if filled_count == 0:
        return "text", {}

    date_matches = sum(
        1 for x in values
        if _RE_DATE1.match(x) or _RE_DATE2.match(x)
    )
    if date_matches / filled_count > 0.4:
        unique_masks = len({_RE_DATE_MASK.sub('X', x) for x in values})
        has_mixed = "yes" if unique_masks > 1 else "no"
        return "date", {
            "inconsistent_date_formatting": has_mixed,
            "desc": "Mixed date formats found." if has_mixed == "yes" else "Dates are uniform."
        }

    at_count = sum(1 for x in values if '@' in x)
    if at_count / filled_count > 0.4:
        has_invalid, has_mixed_case = "no", "no"
        invalid_list = []
        for val in values:
            if not _RE_VALID_EMAIL.match(val):
                has_invalid = "yes"
                if val:
                    invalid_list.append(val)
            if any(c.isupper() for c in val):
                has_mixed_case = "yes"
        return "email", {
            "invalid_emails": has_invalid,
            "invalid_emails_desc": "Column contains broken email formats.",
            "invalid_email_list": list(set(invalid_list)),
            "mixed_case_emails": has_mixed_case,
            "mixed_case_emails_desc": (
                "Emails contain mixed uppercase letters."
                if has_mixed_case == "yes" else "Email casing is uniform."
            )
        }

    phone_matches = sum(
        1 for x in values
        if (d := _RE_NON_DIGIT.sub('', x)).isdigit() and 7 <= len(d) <= 15
    )
    if phone_matches / filled_count > 0.4:
        has_issue, issue_desc = "no", "Phone numbers are uniform."
        for val in values:
            digits_only = re.sub(r'\D', '', val)
            if '?' in val or val.isalpha():
                has_issue = "yes"
                issue_desc = "Phone numbers contain invalid placeholder text symbols (like '??')."
                break
            elif val.startswith(('1','2','3','4','5','6','7','8','9')) and not val.startswith('+'):
                has_issue = "yes"
                issue_desc = "Phone numbers are truncated, missing leading zeros."
                break
            elif digits_only and len(digits_only) < 9:
                has_issue = "yes"
                issue_desc = "Phone numbers contain broken, short sequences missing digits."
                break
        return "phone", {"missing_leading_zeros": has_issue, "desc": issue_desc}

    number_score, has_contamination, decimal_lengths = 0, "no", []
    for val in values:
        cleaned = _RE_NUMBER_CLEAN.sub('', val.replace(',', '.'))
        if re.match(r'^-?\d+(\.\d+)?$', cleaned):
            number_score += 1
            if _RE_CONTAMINATE.search(val) or ('-' in val and not val.strip().startswith('-')):
                has_contamination = "yes"
            decimal_lengths.append(len(cleaned.split('.')[-1]) if '.' in cleaned else 0)
        elif val.lower() in _TEXT_NUMBERS:
            number_score += 1
            has_contamination = "yes"

    if number_score / filled_count > 0.4:
        inconsistent_decimals = "yes" if len(set(decimal_lengths)) > 1 else "no"
        return "number", {
            "inconsistent_numbering": has_contamination,
            "numbering_desc": (
                "Math cells contain currency symbols or text characters."
                if has_contamination == "yes" else "Numbers are cleanly formatted."
            ),
            "inconsistent_decimal_places": inconsistent_decimals,
            "decimal_desc": (
                "Uneven decimal lengths found across numbers."
                if inconsistent_decimals == "yes" else "Decimal lengths are uniform."
            )
        }

    lower_count = upper_count = title_count = 0
    has_newlines = "no"
    for val in values:
        if "\n" in val or "\r" in val:
            has_newlines = "yes"
        if val.islower():
            lower_count += 1
        elif val.isupper():
            upper_count += 1
        elif _RE_TITLE.match(val):
            title_count += 1

    inconsistent_casing = (
        "no" if (lower_count == filled_count or
                 upper_count == filled_count or
                 title_count == filled_count)
        else "yes"
    )
    return "text", {
        "inconsistent_formatting": inconsistent_casing,
        "casing_desc": (
            "Hidden newline breaks found breaking cell format limits."
            if has_newlines == "yes" else "Inconsistent text casing layouts found."
        )
    }


# ─────────────────────────────────────────────
# SUGGEST CLEAN VALUE
# ─────────────────────────────────────────────
def suggest_clean_column(raw_series: pd.Series, detected_type: str,
                         type_metrics: dict, spell_map: dict) -> pd.Series:
    s = raw_series.fillna("").astype(str).str.strip()

    if detected_type == "text":
        def clean_text(val):
            if not val:
                return ""
            words = val.split()
            corrected = []
            for w in words:
                key = _RE_NON_ALPHA_KEY.sub('', w).lower()
                if key in spell_map and spell_map[key]:
                    sg = spell_map[key]
                    if w.isupper():   sg = sg.upper()
                    elif w.istitle(): sg = sg.title()
                    corrected.append(sg)
                else:
                    corrected.append(w)
            result = " ".join(corrected)
            if type_metrics.get("inconsistent_formatting") == "yes":
                result = result.title()
            return result
        return s.apply(clean_text)

    if detected_type == "number":
        def clean_number(val):
            if not val:
                return ""
            cleaned = _RE_NON_DIGIT_SEP.sub('', val.replace(',', '.'))
            try:
                return f"{float(cleaned):.2f}"
            except ValueError:
                return val
        return s.apply(clean_number)

    if detected_type == "date":
        def clean_date(val):
            if not val:
                return ""
            for fmt in _DATE_FORMATS:
                try:
                    return datetime.datetime.strptime(val, fmt).strftime('%Y/%m/%d')
                except ValueError:
                    continue
            return val
        return s.apply(clean_date)

    if detected_type == "email":
        return s.str.lower()

    if detected_type == "phone":
        def clean_phone(val):
            if not val:
                return ""
            digits = re.sub(r'\D', '', val)
            if digits and not val.startswith('+') and len(digits) >= 9:
                return '0' + digits if not digits.startswith('0') else digits
            return val
        return s.apply(clean_phone)

    return s


# ─────────────────────────────────────────────
# EXISTING ROUTE  — untouched
# ─────────────────────────────────────────────
@app.route("/parse-excel", methods=["POST"])
def parse_excel():
    try:
        if "file" not in request.files:
            return jsonify({"error": "No file chunk found"}), 400
        file = request.files["file"]
        if file.filename == "":
            return jsonify({"error": "Empty name"}), 400

        fname = file.filename.lower()
        raw = file.read()

        if fname.endswith('.csv'):
            df = pd.read_csv(io.BytesIO(raw), dtype=str)
        elif fname.endswith('.tsv'):
            df = pd.read_csv(io.BytesIO(raw), sep='\t', dtype=str)
        elif fname.endswith('.ods'):
            df = pd.read_excel(io.BytesIO(raw), dtype=str, engine='odf')
        else:
            df = pd.read_excel(io.BytesIO(raw), dtype=str)

        df.columns = [
            str(col).strip() if not str(col).startswith("Unnamed:") else f"Column {i+1}"
            for i, col in enumerate(df.columns)
        ]

        headers    = list(df.columns)
        total_rows = len(df)

        duplicate_mask    = df.duplicated(keep='first')
        duplicate_indices = [int(i + 1) for i, v in enumerate(duplicate_mask) if v]

        column_diagnostics = {}
        layout_shifts      = []
        spell_maps         = {}
        col_types          = {}
        col_metrics        = {}

        for i, col in enumerate(headers):
            series   = df[col]
            col_vals = (
                series.dropna()
                      .astype(str)
                      .str.strip()
                      .loc[lambda s: s != ""]
                      .tolist()
            )
            filled_count = len(col_vals)

            if total_rows > 5 and filled_count > 0 and (filled_count / total_rows) < 0.15:
                layout_shifts.append({
                    "column": col,
                    "error_msg": f"Stray text boundary shift in {col}.",
                    "sample_value": col_vals[0] if col_vals else ""
                })

            blank_count = int(
                series.isna().sum() +
                (series.astype(str).str.strip() == "").sum()
            )

            detected_type, type_metrics = analyze_exclusive_type(col_vals)
            col_types[col]   = detected_type
            col_metrics[col] = type_metrics

            typos     = []
            spell_map = {}
            if detected_type == "text" and col_vals:
                typos     = find_column_typos(pd.Series(col_vals))
                spell_map = {t["original"]: t["suggestion"] for t in typos}
            spell_maps[col] = spell_map

            mistakes_found = {
                "blank_cells":      blank_count,
                "blank_cells_desc": (
                    f"Found {blank_count} empty rows missing values."
                    if blank_count > 0 else "No missing values."
                )
            }

            if detected_type == "number":
                mistakes_found["inconsistent_numbering"]           = type_metrics.get("inconsistent_numbering", "no")
                mistakes_found["inconsistent_numbering_desc"]      = type_metrics.get("numbering_desc", "")
                mistakes_found["inconsistent_decimal_places"]      = type_metrics.get("inconsistent_decimal_places", "no")
                mistakes_found["inconsistent_decimal_places_desc"] = type_metrics.get("decimal_desc", "")
            elif detected_type == "date":
                mistakes_found["inconsistent_dates_formatting"]      = type_metrics.get("inconsistent_date_formatting", "no")
                mistakes_found["inconsistent_dates_formatting_desc"] = type_metrics.get("desc", "")
            elif detected_type == "phone":
                mistakes_found["missing_leading_zeros"]      = type_metrics.get("missing_leading_zeros", "no")
                mistakes_found["missing_leading_zeros_desc"] = type_metrics.get("desc", "")
            elif detected_type == "email":
                mistakes_found["invalid_emails"]         = type_metrics.get("invalid_emails", "no")
                mistakes_found["invalid_emails_desc"]    = type_metrics.get("invalid_emails_desc", "")
                mistakes_found["invalid_email_list"]     = type_metrics.get("invalid_email_list", [])
                mistakes_found["mixed_case_emails"]      = type_metrics.get("mixed_case_emails", "no")
                mistakes_found["mixed_case_emails_desc"] = type_metrics.get("mixed_case_emails_desc", "")
            elif detected_type == "text":
                mistakes_found["inconsistent_formatting"]      = type_metrics.get("inconsistent_formatting", "no")
                mistakes_found["inconsistent_formatting_desc"] = type_metrics.get("casing_desc", "")
                mistakes_found["misspellings"]                 = typos

            column_diagnostics[col] = {
                "class":          detected_type,
                "mistakes_found": mistakes_found
            }

        cleaned_cols = {}
        for col in headers:
            cleaned_cols[col] = suggest_clean_column(
                df[col],
                col_types[col],
                col_metrics[col],
                spell_maps[col]
            ).tolist()

        suggested_rows = [
            {col: cleaned_cols[col][ri] for col in headers}
            for ri in range(total_rows)
        ]

        df_cleaned = df.fillna("")
        clean_rows = [
            "|".join(str(v) for v in row)
            for row in df_cleaned.itertuples(index=False, name=None)
        ]

        payload = {
            "headers":                 headers,
            "rows_json":               clean_rows,
            "suggested_rows":          suggested_rows,
            "layout_alignment_errors": layout_shifts,
            "duplicate_row_indices":   duplicate_indices,
            "diagnostics":             column_diagnostics
        }

        return jsonify({**payload, "raw_json": json.dumps(payload)}), 200

    except Exception as e:
        return jsonify({"error": f"Internal MasterX workflow crash: {str(e)}"}), 500


# ─────────────────────────────────────────────
# HELPERS FOR /extract-to-excel
# ─────────────────────────────────────────────

def extract_text_from_pdf(raw: bytes) -> str:
    """Extract all text from a PDF using pdfplumber."""
    if not PDF_AVAILABLE:
        raise RuntimeError("pdfplumber is not installed.")
    text_parts = []
    with pdfplumber.open(io.BytesIO(raw)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
    return "\n".join(text_parts)


def extract_text_from_docx(raw: bytes) -> str:
    """Extract all text from a Word .docx file."""
    if not DOCX_AVAILABLE:
        raise RuntimeError("python-docx is not installed.")
    doc = DocxDocument(io.BytesIO(raw))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also pull text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    return "\n".join(paragraphs)


def extract_text_from_pptx(raw: bytes) -> str:
    """Extract all text from a PowerPoint .pptx file."""
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is not installed.")
    prs = Presentation(io.BytesIO(raw))
    lines = []
    for slide_num, slide in enumerate(prs.slides, 1):
        lines.append(f"[Slide {slide_num}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines)


def extract_text_from_image(raw: bytes) -> str:
    """Use OCR (pytesseract) to pull text from an image."""
    if not OCR_AVAILABLE:
        raise RuntimeError("Pillow and pytesseract are not installed.")
    img = Image.open(io.BytesIO(raw))
    return pytesseract.image_to_string(img)


def extract_text_from_txt(raw: bytes) -> str:
    """Decode a plain text file."""
    for enc in ("utf-8", "latin-1", "cp1252"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def extract_text_from_spreadsheet(raw: bytes, fname: str) -> str:
    """Read an existing spreadsheet and return it as tab-separated text."""
    if fname.endswith('.csv'):
        df = pd.read_csv(io.BytesIO(raw), dtype=str)
    elif fname.endswith('.tsv'):
        df = pd.read_csv(io.BytesIO(raw), sep='\t', dtype=str)
    elif fname.endswith('.ods'):
        df = pd.read_excel(io.BytesIO(raw), dtype=str, engine='odf')
    else:
        df = pd.read_excel(io.BytesIO(raw), dtype=str)
    return df.to_csv(index=False)


def parse_text_to_dataframe(text: str) -> pd.DataFrame:
    """
    Intelligently turn raw extracted text into a structured DataFrame.

    Strategy (in order):
    1. If text looks like CSV/TSV already, parse it directly.
    2. Try to detect a table-like structure (consistent delimiters per line).
    3. Fall back: split each non-empty line into a single-column 'Content' table.
    """
    lines = [l for l in text.splitlines() if l.strip()]
    if not lines:
        return pd.DataFrame({"Content": ["(No text could be extracted from this file.)"]})

    # ── Try CSV/TSV detection ──
    first_line = lines[0]
    if '\t' in first_line and first_line.count('\t') >= 1:
        try:
            df = pd.read_csv(io.StringIO(text), sep='\t', dtype=str)
            if len(df.columns) > 1:
                return df.fillna("")
        except Exception:
            pass

    if ',' in first_line and first_line.count(',') >= 1:
        try:
            df = pd.read_csv(io.StringIO(text), dtype=str)
            if len(df.columns) > 1:
                return df.fillna("")
        except Exception:
            pass

    # ── Try to detect a consistent column structure ──
    # Look for lines where words are separated by 2+ spaces (common in PDF tables)
    col_counts = []
    for line in lines[:20]:  # sample first 20 lines
        parts = re.split(r'\s{2,}', line.strip())
        col_counts.append(len(parts))

    if col_counts:
        most_common_cols = max(set(col_counts), key=col_counts.count)
        consistent = sum(1 for c in col_counts if c == most_common_cols)
        # If more than half the lines agree on column count, treat as table
        if most_common_cols > 1 and consistent / len(col_counts) > 0.5:
            rows = []
            for line in lines:
                parts = re.split(r'\s{2,}', line.strip())
                # Pad or trim to match column count
                while len(parts) < most_common_cols:
                    parts.append("")
                rows.append(parts[:most_common_cols])

            if rows:
                header = rows[0]
                data   = rows[1:]
                # If header looks like real column names (not just numbers/dates), use it
                if all(not re.match(r'^\d+$', h) for h in header):
                    df = pd.DataFrame(data, columns=header)
                else:
                    df = pd.DataFrame(rows, columns=[f"Column {i+1}" for i in range(most_common_cols)])
                return df.fillna("")

    # ── Fallback: one row per line, single column ──
    return pd.DataFrame({"Content": lines})


# ─────────────────────────────────────────────
# NEW ROUTE  — independent from /parse-excel
# ─────────────────────────────────────────────
@app.route("/extract-to-excel", methods=["POST"])
def extract_to_excel():
    """
    Accepts any file type (PDF, DOCX, PPTX, image, TXT, CSV, XLSX, etc.),
    extracts the text/data content, structures it into a DataFrame,
    and returns a clean .xlsx file as a download.
    """
    try:
        if "file" not in request.files:
            return jsonify({"error": "No file provided. Send a file with key 'file'."}), 400

        file = request.files["file"]
        if not file.filename:
            return jsonify({"error": "File has no name."}), 400

        fname = file.filename.lower()
        raw   = file.read()

        # ── Route to the right extractor based on file extension ──
        if fname.endswith('.pdf'):
            text = extract_text_from_pdf(raw)
            df   = parse_text_to_dataframe(text)

        elif fname.endswith('.docx'):
            text = extract_text_from_docx(raw)
            df   = parse_text_to_dataframe(text)

        elif fname.endswith('.doc'):
            return jsonify({"error": "Legacy .doc files are not supported. Please convert to .docx first."}), 415

        elif fname.endswith('.pptx'):
            text = extract_text_from_pptx(raw)
            df   = parse_text_to_dataframe(text)

        elif fname.endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif', '.webp', '.gif')):
            text = extract_text_from_image(raw)
            df   = parse_text_to_dataframe(text)

        elif fname.endswith('.txt') or fname.endswith('.md') or fname.endswith('.log'):
            text = extract_text_from_txt(raw)
            df   = parse_text_to_dataframe(text)

        elif fname.endswith(('.xlsx', '.xls', '.csv', '.tsv', '.ods')):
            # Already structured — just clean it up
            text = extract_text_from_spreadsheet(raw, fname)
            df   = parse_text_to_dataframe(text)

        elif fname.endswith('.json'):
            try:
                data = json.loads(raw.decode('utf-8'))
                if isinstance(data, list):
                    df = pd.DataFrame(data).fillna("")
                elif isinstance(data, dict):
                    df = pd.DataFrame([data]).fillna("")
                else:
                    df = pd.DataFrame({"Content": [str(data)]})
            except Exception as e:
                return jsonify({"error": f"Could not parse JSON: {str(e)}"}), 400

        else:
            # Unknown type — try plain text decode as last resort
            try:
                text = extract_text_from_txt(raw)
                df   = parse_text_to_dataframe(text)
            except Exception:
                return jsonify({
                    "error": f"File type '{fname.split('.')[-1]}' is not supported."
                }), 415

        # ── Guard: empty result ──
        if df.empty or (len(df.columns) == 0):
            return jsonify({"error": "Could not extract any structured content from this file."}), 422

        # ── Write to Excel in memory and return as download ──
        output = io.BytesIO()
        with pd.ExcelWriter(output, engine='openpyxl') as writer:
            df.to_excel(writer, index=False, sheet_name='Extracted Data')

            # Auto-size columns for readability
            worksheet = writer.sheets['Extracted Data']
            for col_cells in worksheet.columns:
                max_len = 0
                col_letter = col_cells[0].column_letter
                for cell in col_cells:
                    try:
                        cell_len = len(str(cell.value)) if cell.value else 0
                        if cell_len > max_len:
                            max_len = cell_len
                    except Exception:
                        pass
                adjusted_width = min(max_len + 4, 60)
                worksheet.column_dimensions[col_letter].width = adjusted_width

        output.seek(0)

        # ── Build a clean output filename ──
        base_name = re.sub(r'\.[^.]+$', '', file.filename)  # strip extension
        safe_name = re.sub(r'[^\w\-]', '_', base_name)
        download_name = f"{safe_name}_extracted.xlsx"

        return send_file(
            output,
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            as_attachment=True,
            download_name=download_name
        )

    except RuntimeError as e:
        return jsonify({"error": f"Missing library: {str(e)}"}), 500
    except Exception as e:
        return jsonify({"error": f"Extraction failed: {str(e)}"}), 500


# ─────────────────────────────────────────────
# ENTRY POINT
# ─────────────────────────────────────────────
if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port)
